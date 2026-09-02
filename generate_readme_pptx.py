#!/usr/bin/env python3
"""
Generate a Cisco CX-branded PPTX presentation from README.md content.
Usage: python generate_readme_pptx.py
Output: README-presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy
import os

# ── Paths ────────────────────────────────────────────────────────────────────
REPO = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = os.path.join(REPO, "backend/resources/pptAssets/cxPpt_template.pptx")
ARCH_IMG = os.path.join(REPO, "architecture_diagram_advanced_health_check.png")
OUTPUT   = os.path.join(REPO, "README-presentation.pptx")

# ── Cisco brand colors ────────────────────────────────────────────────────────
CISCO_BLUE   = RGBColor(0x00, 0x52, 0x73)   # #005273 dark blue
CISCO_TEAL   = RGBColor(0x00, 0x9F, 0xCC)   # #009FCC light blue
CISCO_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CISCO_DARK   = RGBColor(0x1A, 0x1A, 0x1A)
CISCO_GRAY   = RGBColor(0x58, 0x58, 0x58)
CISCO_LIGHT  = RGBColor(0xF4, 0xF4, 0xF4)
ACCENT_AMBER = RGBColor(0xFF, 0xA5, 0x00)

W = Inches(13.33)  # slide width
H = Inches(7.50)   # slide height

# ── Helpers ───────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    fill_fmt = shape.fill
    if fill:
        fill_fmt.solid()
        fill_fmt.fore_color.rgb = fill
    else:
        fill_fmt.background()
    line_fmt = shape.line
    if line:
        line_fmt.color.rgb = line
    else:
        line_fmt.fill.background()
    return shape


def add_textbox(slide, l, t, w, h, text, size=18, bold=False, color=CISCO_DARK,
                align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def add_bullet_textbox(slide, l, t, w, h, items, size=16, color=CISCO_DARK,
                       title=None, title_size=22, title_color=CISCO_BLUE,
                       bullet_char="•  ", indent=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(title_size)
        run.font.bold = True
        run.font.color.rgb = title_color
    for i, item in enumerate(items):
        if title or i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        if indent:
            p.level = 1
        run = p.add_run()
        run.text = f"{bullet_char}{item}"
        run.font.size  = Pt(size)
        run.font.color.rgb = color
        run.font.bold  = False
    return txb


def slide_chrome(slide, title_text, bar_color=CISCO_BLUE, text_color=CISCO_WHITE,
                 subtitle=None):
    """Add a top bar + title text to a slide."""
    # top accent bar
    add_rect(slide, 0, 0, W, Inches(0.07), fill=CISCO_TEAL)
    # title bar
    add_rect(slide, 0, Inches(0.07), W, Inches(0.85), fill=bar_color)
    # title text
    add_textbox(slide, Inches(0.35), Inches(0.10), W - Inches(0.7), Inches(0.82),
                title_text, size=24, bold=True, color=text_color)
    # optional subtitle
    if subtitle:
        add_textbox(slide, Inches(0.35), Inches(0.60), W - Inches(0.7), Inches(0.40),
                    subtitle, size=13, bold=False, color=CISCO_TEAL)
    # bottom bar
    add_rect(slide, 0, H - Inches(0.30), W, Inches(0.30), fill=CISCO_BLUE)
    add_textbox(slide, Inches(0.2), H - Inches(0.30), Inches(6), Inches(0.30),
                "Cisco CX Hackathon 2026  ·  CS-IC-62  ·  Bridging the Observability Gap",
                size=8, color=CISCO_WHITE)


def add_table_slide(slide, rows, col_widths, l, t, w, h,
                    header_fill=CISCO_BLUE, header_text_color=CISCO_WHITE,
                    row_fill=CISCO_WHITE, alt_fill=CISCO_LIGHT,
                    font_size=10, header_font_size=11):
    from pptx.util import Pt
    n_rows = len(rows)
    n_cols = len(rows[0])
    tbl = slide.shapes.add_table(n_rows, n_cols, l, t, w, h).table

    # column widths
    total_w = sum(col_widths)
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = int(w * cw / total_w)

    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = str(cell_text)
            tf = cell.text_frame
            tf.word_wrap = True
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(header_font_size if ri == 0 else font_size)
                    run.font.bold = (ri == 0)
                    run.font.color.rgb = header_text_color if ri == 0 else CISCO_DARK
            # background
            fill = header_fill if ri == 0 else (alt_fill if ri % 2 == 0 else row_fill)
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
            srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
            srgb.set('val', str(fill))
    return tbl


# ── Presentation builder ──────────────────────────────────────────────────────
def build():
    # Start from a truly blank presentation (no template baggage)
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    # Use the blank layout (index 6 in the default blank deck)
    blank_layout = prs.slide_layouts[6]

    # ── SLIDE 1: Title ─────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    # full background gradient — top dark, bottom slightly lighter
    add_rect(s, 0, 0, W, H, fill=CISCO_BLUE)
    add_rect(s, 0, 0, W, Inches(0.07), fill=CISCO_TEAL)
    add_rect(s, 0, H - Inches(0.07), W, Inches(0.07), fill=CISCO_TEAL)
    # diagonal accent strip
    add_rect(s, W - Inches(3.5), 0, Inches(3.5), H, fill=RGBColor(0x00, 0x42, 0x60))

    add_textbox(s, Inches(0.6), Inches(1.6), Inches(8.5), Inches(1.4),
                "Bridging the Observability Gap", size=38, bold=True,
                color=CISCO_WHITE, align=PP_ALIGN.LEFT)
    add_textbox(s, Inches(0.6), Inches(3.0), Inches(8.5), Inches(0.8),
                "with AI-Driven Insights", size=32, bold=False,
                color=CISCO_TEAL, align=PP_ALIGN.LEFT)
    add_textbox(s, Inches(0.6), Inches(4.0), Inches(8.5), Inches(0.5),
                "Cisco CX Hackathon 2026  ·  Team CS-IC-62",
                size=16, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.LEFT)

    add_textbox(s, Inches(0.6), Inches(5.5), Inches(8), Inches(0.6),
                "MCP  ·  CIRCUIT LLM  ·  AppDynamics CAT  ·  AI-Driven Insights",
                size=13, color=CISCO_TEAL, italic=True)

    # ── SLIDE 2: Agenda ────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    slide_chrome(s, "Agenda")
    topics = [
        "1.  The Problem — why manual assessment doesn't scale",
        "2.  Solution Overview — the AI-driven health check pipeline",
        "3.  Architecture — how MCP, CIRCUIT LLM & CAT work together",
        "4.  LLM Role — what the AI does at each stage",
        "5.  Why LLM is the right approach",
        "6.  Live Demo — run, analyse, chart, Q&A",
        "7.  Key Outcomes — measured results",
        "8.  Setup & Next Steps",
    ]
    add_bullet_textbox(s, Inches(0.8), Inches(1.2), Inches(11.5), Inches(5.6),
                       topics, size=17, bullet_char="", color=CISCO_DARK)

    # ── SLIDE 3: Problem Statement ────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    slide_chrome(s, "The Problem — Manual Assessment Doesn't Scale")

    problems = [
        "Time-to-insight is too slow — multi-tab Excel & PowerPoint require 30–60 min of manual review before any recommendation reaches the customer.",
        "Expertise barrier is high — deep AppDynamics API knowledge and maturity-model familiarity don't scale across the entire CS org.",
        "Reports go stale and get lost — assessment outputs are rarely archived consistently; trend analysis is nearly impossible.",
        "Follow-up questions restart the cycle — every clarifying question ('Which apps are blocking Gold tier?') triggers another manual report review.",
    ]
    numbers = ["01", "02", "03", "04"]
    cols = [0, 1, 0, 1]
    rows_pos = [0, 0, 1, 1]
    col_x = [Inches(0.4), Inches(6.9)]
    row_y = [Inches(1.15), Inches(3.90)]
    bw, bh = Inches(6.1), Inches(2.4)

    for i, (prob, num) in enumerate(zip(problems, numbers)):
        x = col_x[cols[i]]
        y = row_y[rows_pos[i]]
        add_rect(s, x, y, bw, bh, fill=CISCO_LIGHT)
        add_rect(s, x, y, Inches(0.55), bh, fill=CISCO_BLUE)
        add_textbox(s, x + Inches(0.05), y + Inches(0.6), Inches(0.5), Inches(0.8),
                    num, size=20, bold=True, color=CISCO_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.65), y + Inches(0.15), bw - Inches(0.75), bh - Inches(0.3),
                    prob, size=13, color=CISCO_DARK)

    add_textbox(s, Inches(1), H - Inches(0.85), W - Inches(2), Inches(0.5),
                "Result: observability data that exists but doesn't flow. The gap between raw telemetry and business-level guidance stays wide — and expensive.",
                size=11, italic=True, color=CISCO_GRAY)

    # ── SLIDE 4: Solution Overview ────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    slide_chrome(s, "Solution — AI-Driven Health Check Pipeline")

    add_textbox(s, Inches(0.5), Inches(1.05), Inches(12), Inches(0.5),
                "Automates the full assessment lifecycle — from metric collection to report delivery to conversational Q&A — using:",
                size=14, color=CISCO_DARK)

    pillars = [
        ("MCP", "Model Context Protocol", "Structured tool-calling between LLM and backend engine"),
        ("CIRCUIT LLM", "Cisco's internal AI gateway", "Secure, on-network reasoning; no data leaves Cisco"),
        ("CAT Engine", "Config Assessment Tool", "Deterministic AppDynamics API extraction + Excel/PPT reports"),
        ("Knowledge Base", "Persistent context store", "Historical reports feed cumulative cross-run intelligence"),
    ]
    pw = Inches(3.0)
    ph = Inches(3.2)
    gap = Inches(0.13)
    for i, (abbr, name, desc) in enumerate(pillars):
        x = Inches(0.35) + i * (pw + gap)
        y = Inches(1.65)
        add_rect(s, x, y, pw, ph, fill=CISCO_BLUE)
        add_textbox(s, x + Inches(0.12), y + Inches(0.18), pw - Inches(0.24), Inches(0.70),
                    abbr, size=26, bold=True, color=CISCO_TEAL, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.12), y + Inches(0.90), pw - Inches(0.24), Inches(0.5),
                    name, size=13, bold=True, color=CISCO_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.12), y + Inches(1.50), pw - Inches(0.24), Inches(1.5),
                    desc, size=11, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)

    add_textbox(s, Inches(0.5), H - Inches(0.85), W - Inches(1), Inches(0.5),
                "Every run is archived to SharePoint. Report data is injected into LLM context for instant conversational Q&A — no re-running required.",
                size=11, italic=True, color=CISCO_GRAY)

    # ── SLIDE 5: Architecture Diagram ─────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    slide_chrome(s, "Architecture — Advanced Health Check & LLM Insight Flow")
    if os.path.exists(ARCH_IMG):
        # center image below the title bar
        img_w = Inches(12.0)
        img_h = Inches(5.8)
        img_l = (W - img_w) / 2
        img_t = Inches(1.00)
        s.shapes.add_picture(ARCH_IMG, img_l, img_t, img_w, img_h)
    else:
        add_textbox(s, Inches(2), Inches(3), Inches(9), Inches(1),
                    "[architecture_diagram_advanced_health_check.png not found]",
                    size=16, color=CISCO_GRAY, align=PP_ALIGN.CENTER)

    # ── SLIDE 6: How the LLM Powers Each Stage ────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    slide_chrome(s, "How the LLM Powers Each Stage")

    table_data = [
        ["Step", "What Happens", "LLM's Role"],
        ["1 — Request",         "User asks for a health check in plain English",               "Parses intent, identifies target environment & job parameters"],
        ["2 — Orchestration",   "CIRCUIT hands off a structured tool call to the MCP Server",  "Selects the correct MCP tool and arguments autonomously"],
        ["3 — Collection",      "Health Endpoints Engine queries AppDynamics Controller APIs",  "— deterministic engine —"],
        ["4 — Report Gen",      "Engine produces Excel and PowerPoint reports",                 "— deterministic engine —"],
        ["5 — Archiving",       "Archive Agent saves reports to SharePoint Drive",              "LLM can trigger archiving in the same conversation turn"],
        ["6 — Knowledge",       "Report data is scraped and fed into the LLM's context",       "Builds cumulative understanding across runs and customers"],
        ["7 — Delivery",        "Reports attached to CIRCUIT prompt, available for download",   "Narrates key findings and highlights critical gaps"],
        ["8 — Q&A",             "User asks follow-up questions",                               "Answers from loaded context — no re-running any tools"],
    ]
    col_widths = [1.8, 4.5, 4.5]
    add_table_slide(s, table_data, col_widths,
                    l=Inches(0.35), t=Inches(1.0), w=Inches(12.6), h=Inches(6.0),
                    font_size=10, header_font_size=12)

    # ── SLIDE 7: Why LLM ──────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    slide_chrome(s, "Why LLM is the Right Approach")

    reasons = [
        (
            "Tool-Use Reasoning",
            "CIRCUIT receives the full MCP tool schema at session start and autonomously decides which tool to call, when, and with what arguments — no scripting required.",
        ),
        (
            "Structured → Natural Translation",
            "Dense multi-sheet Excel data (maturity scores, agent matrices, threshold breaches) is converted into concise, human-readable summaries and prioritised recommendations.",
        ),
        (
            "Persistent Knowledge Base",
            "Each run feeds the LLM's context, enabling cross-run and cross-customer pattern recognition that no single-session tool can provide.",
        ),
        (
            "Conversational Memory",
            "The full report is injected into context once. Follow-up questions are answered from memory — no re-running, no re-downloading.",
        ),
    ]
    rw, rh = Inches(5.9), Inches(2.4)
    positions = [(Inches(0.35), Inches(1.10)), (Inches(6.65), Inches(1.10)),
                 (Inches(0.35), Inches(3.80)), (Inches(6.65), Inches(3.80))]
    icons = ["🔧", "🔄", "🧠", "💬"]

    for (title, body), (x, y), icon in zip(reasons, positions, icons):
        add_rect(s, x, y, rw, rh, fill=CISCO_LIGHT)
        add_rect(s, x, y, rw, Inches(0.5), fill=CISCO_TEAL)
        add_textbox(s, x + Inches(0.15), y + Inches(0.05), rw - Inches(0.3), Inches(0.45),
                    title, size=14, bold=True, color=CISCO_WHITE)
        add_textbox(s, x + Inches(0.15), y + Inches(0.6), rw - Inches(0.3), rh - Inches(0.7),
                    body, size=12, color=CISCO_DARK)

    # ── SLIDE 8: Demo Flow ────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    slide_chrome(s, "Live Demo — Assessment to Insight in One Conversation")

    steps = [
        ("1", "Run assessment", 'User: "Run health check for DefaultJob"', CISCO_BLUE),
        ("2", "Instant summary", "CIRCUIT: Overall APM maturity is 68% (Silver). 3 apps blocked from Gold.", CISCO_TEAL),
        ("3", "Drill down", 'User: "Which apps are furthest from Gold?"', CISCO_BLUE),
        ("4", "LLM answers", "CIRCUIT answers from loaded context — no tool re-run needed.", CISCO_TEAL),
        ("5", "Archive", 'User: "Archive this and send the download link"', CISCO_BLUE),
        ("6", "Done", "CIRCUIT: Report archived to SharePoint. Download link: https://...", CISCO_TEAL),
    ]
    sw = Inches(3.9)
    sh = Inches(2.15)
    gap = Inches(0.15)
    for i, (num, title, desc, color) in enumerate(steps):
        col = i % 3
        row = i // 3
        x = Inches(0.35) + col * (sw + gap)
        y = Inches(1.15) + row * (sh + gap)
        add_rect(s, x, y, sw, sh, fill=color)
        add_textbox(s, x + Inches(0.1), y + Inches(0.08), Inches(0.55), Inches(0.5),
                    num, size=22, bold=True, color=CISCO_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.08), y + Inches(0.55), sw - Inches(0.16), Inches(0.45),
                    title, size=13, bold=True, color=CISCO_WHITE)
        add_textbox(s, x + Inches(0.08), y + Inches(1.05), sw - Inches(0.16), sh - Inches(1.15),
                    desc, size=10, color=RGBColor(0xCC, 0xEE, 0xFF) if color == CISCO_BLUE else CISCO_WHITE)

    add_textbox(s, Inches(0.5), H - Inches(0.85), W - Inches(1), Inches(0.4),
                "Tip: the LLM answers follow-up questions from in-context memory — no re-running, no re-downloading reports.",
                size=10, italic=True, color=CISCO_GRAY)

    # ── SLIDE 9: Key Outcomes ──────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    slide_chrome(s, "Key Outcomes — Measured Results")

    outcomes = [
        ("10×", "Faster Time-to-Insight",     "Results previously requiring 30–60 min of manual Excel review\nsummarised in seconds"),
        ("0",   "Tool Expertise Required",    "Any CS engineer or SA can request a health check and receive\nconsultant-grade guidance via natural language"),
        ("100%","Automated Archiving",         "Every run saved to SharePoint via Archive Agent —\nsearchable, time-stamped, zero manual effort"),
        ("∞",   "Cumulative Knowledge",        "Report data feeds the LLM knowledge base enabling\ncross-run trend & cross-customer pattern recognition"),
        ("1×",  "Context Load, Any Questions", "The full report is injected into context once;\nall follow-up Q&A answered from memory"),
        ("✓",   "Cisco-Native & Secure",       "Runs entirely over CIRCUIT (Cisco's internal AI gateway) —\nall data stays on-network"),
    ]
    ow = Inches(4.1)
    oh = Inches(2.05)
    gap = Inches(0.10)
    for i, (stat, title, desc) in enumerate(outcomes):
        col = i % 3
        row = i // 3
        x = Inches(0.35) + col * (ow + gap)
        y = Inches(1.10) + row * (oh + gap)
        add_rect(s, x, y, ow, oh, fill=CISCO_LIGHT)
        add_rect(s, x, y, Inches(1.2), oh, fill=CISCO_BLUE)
        add_textbox(s, x + Inches(0.05), y + Inches(0.5), Inches(1.1), Inches(0.8),
                    stat, size=22, bold=True, color=CISCO_TEAL, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(1.28), y + Inches(0.08), ow - Inches(1.38), Inches(0.45),
                    title, size=13, bold=True, color=CISCO_BLUE)
        add_textbox(s, x + Inches(1.28), y + Inches(0.55), ow - Inches(1.38), oh - Inches(0.65),
                    desc, size=10, color=CISCO_DARK)

    # ── SLIDE 10: Setup & Quick Start ─────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    slide_chrome(s, "Setup & Quick Start")

    prereqs = [
        "Python 3.12+",
        "Cisco CIRCUIT LLM API credentials (app key, client ID, client secret)",
        "AppDynamics controller + a CAT job file in input/jobs/",
    ]
    steps_list = [
        "git clone <repo> && cd <repo>",
        "pip install pipenv && pipenv install && pipenv shell",
        "Copy properties.py.template → properties.py  and fill in CIRCUIT credentials",
        "Create input/jobs/DefaultJob.json with your controller URL + credentials",
        "python mcp_llm_client.py  — and start chatting!",
    ]

    add_textbox(s, Inches(0.5), Inches(1.05), Inches(5.5), Inches(0.4),
                "Prerequisites", size=16, bold=True, color=CISCO_BLUE)
    add_bullet_textbox(s, Inches(0.5), Inches(1.50), Inches(5.5), Inches(1.8),
                       prereqs, size=13, color=CISCO_DARK, bullet_char="✓  ")

    add_textbox(s, Inches(6.5), Inches(1.05), Inches(6.4), Inches(0.4),
                "5 Steps to First Insight", size=16, bold=True, color=CISCO_BLUE)
    for i, step in enumerate(steps_list):
        x, y = Inches(6.5), Inches(1.50) + i * Inches(0.72)
        add_rect(s, x, y, Inches(0.5), Inches(0.5), fill=CISCO_BLUE)
        add_textbox(s, x, y + Inches(0.02), Inches(0.5), Inches(0.46),
                    str(i + 1), size=16, bold=True, color=CISCO_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.6), y + Inches(0.05), Inches(5.7), Inches(0.56),
                    step, size=12, color=CISCO_DARK)

    add_textbox(s, Inches(0.5), Inches(3.5), Inches(5.5), Inches(0.4),
                "Key Environment Variables", size=14, bold=True, color=CISCO_BLUE)
    env_vars = [
        "CIRCUIT_LLM_API_APP_KEY",
        "CIRCUIT_LLM_API_CLIENT_ID",
        "CIRCUIT_LLM_API_CLIENT_SECRET",
        "CIRCUIT_LLM_API_MODEL_NAME   (e.g. gpt-5-nano)",
    ]
    add_bullet_textbox(s, Inches(0.5), Inches(3.95), Inches(5.5), Inches(2.5),
                       env_vars, size=12, color=CISCO_DARK, bullet_char="$  ")

    # ── SLIDE 11: Thank You / Q&A ──────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    add_rect(s, 0, 0, W, H, fill=CISCO_BLUE)
    add_rect(s, 0, 0, W, Inches(0.07), fill=CISCO_TEAL)
    add_rect(s, 0, H - Inches(0.07), W, Inches(0.07), fill=CISCO_TEAL)
    add_rect(s, W - Inches(3.5), 0, Inches(3.5), H, fill=RGBColor(0x00, 0x42, 0x60))

    add_textbox(s, Inches(0.8), Inches(2.2), Inches(8), Inches(1.2),
                "Thank You", size=48, bold=True, color=CISCO_WHITE)
    add_textbox(s, Inches(0.8), Inches(3.4), Inches(8), Inches(0.6),
                "Questions?", size=28, color=CISCO_TEAL)
    add_textbox(s, Inches(0.8), Inches(4.4), Inches(8), Inches(0.5),
                "Team CS-IC-62  ·  Cisco CX Hackathon 2026",
                size=16, color=RGBColor(0xCC, 0xDD, 0xEE))
    add_textbox(s, Inches(0.8), Inches(5.2), Inches(8), Inches(0.5),
                "github: alexafshar/config-assessment-tool  (branch: feature/experiment/hack)",
                size=12, color=RGBColor(0x88, 0xBB, 0xDD), italic=True)

    # ── Save ──────────────────────────────────────────────────────────────────
    prs.save(OUTPUT)
    size_kb = os.path.getsize(OUTPUT) // 1024
    print(f"Saved: {OUTPUT}  ({size_kb} KB, {len(prs.slides)} slides)")


if __name__ == "__main__":
    build()

